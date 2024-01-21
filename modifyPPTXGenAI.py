import os
import sys
import shutil
import zipfile
import hashlib
import openai
import requests
import base64
import io
import subprocess
from pptx import Presentation
from pptx.util import Pt, Inches
import xml.etree.ElementTree as ET
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from collections import defaultdict
from PyPDF2 import PdfFileReader, PdfFileWriter



def ascii_to_braille(text):
    """
    Converts ASCII text to Braille using the 'louis' library.

    Parameters:
    text (str): The ASCII text to convert.

    Returns:
    str: The Braille equivalent of the input text.
    """
    non_translatable_chars = ['▪', ' ', 'A.', 'B.', 'C.', 'D.', 'E.']
    if text.strip() and text.strip() in non_translatable_chars:
        return text

    process = subprocess.Popen(['lou_translate', 'en-ueb-g2.ctb'],
                               stdin=subprocess.PIPE,
                               stdout=subprocess.PIPE,
                               stderr=subprocess.PIPE)
    stdout, stderr = process.communicate(input=text.encode())
    if process.returncode != 0:
        raise Exception(f"lou_translate failed: {stderr.decode()}")
    return stdout.decode()


def hash_file(filepath):
    """Return the SHA-1 hash of the file."""
    BUF_SIZE = 65536  # Read in 64kb chunks
    sha1 = hashlib.sha1()

    with open(filepath, 'rb') as f:
        while True:
            data = f.read(BUF_SIZE)
            if not data:
                break
            sha1.update(data)
    return sha1.hexdigest()
    

def modify_presentation_font(presentation, font_name, font_size, spacing_size=50):
    """
    Modifies the font name and size for all text in a PowerPoint presentation.

    Parameters:
    presentation (pptx.Presentation): The presentation to modify.
    font_name (str): The name of the font to use.
    font_size (int): The size of the font.

    Returns:
    pptx.Presentation: The modified presentation.
    """
    answer = ''
    while answer not in ['yes', 'no']:
        answer = input(f"Do you want to change the font size to what's specified in this script ({font_size}) (yes) or keep them the same (no)? (yes/no):")
        if answer not in ['yes', 'no']:
            print(f"Your answer {answer} is not 'yes' or 'no', please enter one of them")

    change_font_size = (answer == 'yes')
            
    for slide in presentation.slides:
        shapes_to_process = list(slide.shapes)
        while shapes_to_process:
            shape = shapes_to_process.pop()
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.space = Pt(spacing_size)
                        run.font.italic = False
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        if change_font_size:
                            run.font.size = Pt(font_size)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = font_name
                                run.font.space = Pt(spacing_size)
                                run.font.italic = False
                                run.font.color.rgb = RGBColor(0, 0, 0)
                                if change_font_size:
                                    run.font.size = Pt(font_size)
            elif shape.shape_type == 6:  # Group shape
                shapes_to_process.extend(shape.shapes)
    return presentation


def modify_presentation_spacing(presentation, line_spacing_value):
    """
    Modifies the line spacing for all text in a PowerPoint presentation.

    Parameters:
    presentation (pptx.Presentation): The presentation to modify.
    line_spacing_value (float): The line spacing value to set.

    Returns:
    pptx.Presentation: The modified presentation.
    """
    for slide in presentation.slides:
        shapes_to_process = list(slide.shapes)
        while shapes_to_process:
            shape = shapes_to_process.pop()
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.line_spacing = line_spacing_value
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.line_spacing = line_spacing_value
            elif shape.shape_type == 6:  # Group shape
                shapes_to_process.extend(shape.shapes)
    return presentation


def contract_braille(presentation, verbose=False):
    """
    Applies Braille contraction to all text in a PowerPoint presentation.

    Parameters:
    presentation (pptx.Presentation): The presentation to modify.
    verbose (bool): Flag to enable verbose output.

    Returns:
    pptx.Presentation: The modified presentation.
    """
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if verbose:
                            print(f'Original: {run.text}')
                        run.text = ascii_to_braille(run.text)
                        if verbose:
                            print(f'Modified: {run.text}')
    return presentation


def remove_repetitive_images(pptx_path, image_dir):
    """
    Processes files in a directory based on the slide count of a PowerPoint presentation.
    It hashes all files in the directory, and if a hash repeats for more than 80% of the
    slide count in the PowerPoint file, all files with that hash are deleted.

    Parameters:
    pptx_path (str): Path to the PowerPoint file.
    image_dir (str): Path to the directory containing images.

    Returns:
    None
    """

    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)

    hash_counts = defaultdict(int)
    file_hashes = defaultdict(list)

    for filename in os.listdir(image_dir):
        file_path = os.path.join(image_dir, filename)
        if os.path.isfile(file_path):
            file_hash = hash_file(file_path)
            hash_counts[file_hash] += 1
            file_hashes[file_hash].append(file_path)
    print(hash_counts)  

    threshold = 0.8 * num_slides
    for file_hash, files in file_hashes.items():
        if hash_counts[file_hash] > threshold:
            for file_path in files:
                os.remove(file_path)
                print(f"Deleted {file_path}")
                
                
def clean_pptx(pptx_path, output_path):
    """
    Cleans a PowerPoint (PPTX) file by removing references to non-existing images.

    This function decompresses the PPTX file, parses the XML structure to find and 
    remove references to missing images, and then repacks the cleaned presentation 
    into a new PPTX file.

    Args:
    pptx_path (str): The file path of the input PowerPoint file.
    output_path (str): The file path where the cleaned PowerPoint file will be saved.

    Returns:
    None: The function does not return a value but saves the cleaned PPTX file at output_path.
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"The specified PPTX file does not exist: {pptx_path}")

    temp_dir = 'temp_pptx'
    os.makedirs(temp_dir, exist_ok=True)

    with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
        zip_ref.extractall(temp_dir)
        
    image_dir = os.path.join('temp_pptx', os.path.join('ppt', 'media'))
    remove_repetitive_images(pptx_path, image_dir)

    namespaces = {
        'r': 'http://schemas.openxmlformats.org/package/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
    }

    for rels_file in os.listdir(os.path.join(temp_dir, 'ppt', 'slides', '_rels')):
        if not rels_file.endswith('.rels'):
            continue

        tree = ET.parse(os.path.join(temp_dir, 'ppt', 'slides', '_rels', rels_file))
        root = tree.getroot()

        removed_ids = []
        for relationship in root.findall('r:Relationship', namespaces):
            if 'image' in relationship.attrib['Type']:
                image_path = os.path.join(temp_dir, 'ppt', relationship.attrib['Target'].lstrip('/'))
                image_path = str(os.path.join(temp_dir, 'ppt', relationship.attrib['Target'].lstrip('/')).replace('../', ''))
                #print(image_path)
                if not os.path.exists(image_path):
                    removed_ids.append(relationship.attrib['Id'])
                    root.remove(relationship)

        tree.write(os.path.join(temp_dir, 'ppt', 'slides', '_rels', rels_file))

        slide_file = rels_file.replace('.rels', '')
        slide_tree = ET.parse(os.path.join(temp_dir, 'ppt', 'slides', slide_file))
        slide_root = slide_tree.getroot()

        for pic in slide_root.findall('.//p:pic', namespaces):
            blip = pic.find('.//a:blip', namespaces)
            if blip is not None and blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'] in removed_ids:
                pic.getparent().remove(pic)

        slide_tree.write(os.path.join(temp_dir, 'ppt', 'slides', slide_file))

    with zipfile.ZipFile(output_path, 'w') as myzip:
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                myzip.write(os.path.join(root, file), arcname=os.path.relpath(os.path.join(root, file), temp_dir))

    shutil.rmtree(temp_dir)
    
    
def query_gpt4_with_image(image_input, api_key):
    """
    Queries GPT-4 with an image and returns the text output.

    :param image_input: Path to the image file or a PIL image object.
    :param api_key: API key for accessing the GPT-4 service.
    :return: Text response from GPT-4.
    """
    
    def encode_image(image_input):
        if isinstance(image_input, Image.Image):
            if image_input.mode == 'RGBA':
                image_input = image_input.convert('RGB')

            img_byte_arr = io.BytesIO()
            image_input.save(img_byte_arr, format='JPEG')
            return base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
        else:
            with open(image_input, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')

    base64_image = encode_image(image_input)

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }

    instructions = "Describe the image below to a blind person. Do not give definitions of the image's contents.\
    Assume your audience is familiar with the concepts in the image and just needs a description of the image.\
    Make no mention of color.\
    Focus on describing the locations of the things you describe and how they are graphically depicted in the image.\
    Use no more than 200 words."

    payload = {
        "model": "gpt-4-vision-preview",  
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": instructions},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }
        ],
        "max_tokens": 1000
    }

    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)

    if response.status_code == 200:
        return response.json()['choices'][0]['message']['content']
    else:
        raise Exception(f"Error querying GPT-4 with image: {response.status_code} - {response.text}")


    
def create_notes_slides(pptx_file, new_filename):
    """
    Creates a new PowerPoint presentation containing the notes from the original presentation.

    Parameters:
    pptx_file (str): The path to the original PowerPoint file.
    new_filename (str): The path for the new PowerPoint file with notes.
    """
    if not os.path.exists(pptx_file):
        raise FileNotFoundError(f'The specified PPTX file does not exist: {pptx_file}')

    original_presentation = Presentation(pptx_file)
    new_presentation = Presentation()
    slide_width = new_presentation.slide_width
    slide_height = new_presentation.slide_height
    margin = Inches(0.5)

    for index, slide in enumerate(original_presentation.slides):
        notes_slide = slide.notes_slide
        notes_text = 'This string will be replaced with the output of a function that generates detailed text descriptions of images.'
        new_slide_layout = new_presentation.slide_layouts[6]
        notes_slide = new_presentation.slides.add_slide(new_slide_layout)
        text_box = notes_slide.shapes.add_textbox(margin, margin, slide_width - 2 * margin, slide_height - 2 * margin)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = notes_text

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
            paragraph.alignment = PP_ALIGN.LEFT

    new_presentation.save(new_filename)


def extract_images_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    pil_images = {}
    for i, slide in enumerate(prs.slides):
        pil_images[i] = []
        for shape in slide.shapes:
            if shape.shape_type == 13: 
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)

                pil_images[i].append(image)

    return pil_images
    

def convert_pptx_to_pdf(pptx_file):
    """
    Converts a PowerPoint file to a PDF.

    Parameters:
    pptx_file (str): The path to the PowerPoint file to convert.
    """
    output_pdf = os.path.splitext(pptx_file)[0] + '.pdf'

    if not os.path.exists(pptx_file):
        raise FileNotFoundError(f'The specified PPTX file does not exist: {pptx_file}')

    try:
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', pptx_file, '--outdir', os.path.dirname(output_pdf)])
        print(f'File converted successfully. Saved as {output_pdf}.')
    except Exception as error:
        print(f'An error occurred during file conversion: {error}')


def interleave_pdfs(file1, file2, output_file):
    """
    Interleaves pages from two PDF files into a single output file.

    Parameters:
    file1 (str): The path to the first PDF file.
    file2 (str): The path to the second PDF file.
    output_file (str): The path for the output interleaved PDF file.
    """
    if not os.path.exists(file1):
        raise FileNotFoundError(f'The specified PDF file does not exist: {file1}')
    if not os.path.exists(file2):
        raise FileNotFoundError(f'The specified PDF file does not exist: {file2}')

    try:
        pdf1 = PdfFileReader(file1)
        pdf2 = PdfFileReader(file2)
        pdf_writer = PdfFileWriter()
        for i in range(max(pdf1.getNumPages(), pdf2.getNumPages())):
            if i < pdf1.getNumPages():
                pdf_writer.addPage(pdf1.getPage(i))
            if i < pdf2.getNumPages():
                pdf_writer.addPage(pdf2.getPage(i))

        with open(output_file, 'wb') as output:
            pdf_writer.write(output)

        print(f'PDFs interleaved successfully and saved to {output_file}.')
    except Exception as error:
        print(f'An error occurred during PDF interleaving: {error}')

def main():   
    # Clean unncessary images
    new_filename = sys.argv[1].rsplit('.', 1)[0] + '_cleaned_images.pptx'
    clean_pptx(sys.argv[1], new_filename)
    
    # Create 'notes' presentation
    notes_filename = sys.argv[1].rsplit('.', 1)[0] + '_notes.pptx'
    print(notes_filename)
    create_notes_slides(new_filename, notes_filename)

    # Load the lecture slides
    presentation = Presentation(new_filename)
    #presentation = modify_presentation_font(presentation=presentation, font_name='Braille', font_size=14)
    presentation = modify_presentation_spacing(presentation, 1.2)
    #presentation = contract_braille(presentation)

    # Save the modified presentation
    base = os.path.splitext(sys.argv[1])[0]
    new_filename = base + "_braille.pptx"
    presentation.save(new_filename)
    convert_pptx_to_pdf(new_filename)

    # Load the notes presentation
    notes_presentation = Presentation(notes_filename)
    #notes_presentation = modify_presentation_font(presentation=notes_presentation, font_name='Braille', font_size=14)
    #notes_presentation = contract_braille(notes_presentation)
    notes_presentation.save(notes_filename)
    convert_pptx_to_pdf(notes_filename)
    
    # Combine the lecture slides and 'notes' slides
    interleaved_filename = os.path.splitext(new_filename)[0] + '_with_notes.pdf'
    interleave_pdfs(new_filename[:-5] + '.pdf', notes_filename[:-5] + '.pdf', interleaved_filename)
   
if __name__ == '__main__':
    main()
